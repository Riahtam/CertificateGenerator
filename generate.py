from pptx import Presentation
import pandas as pd
from pptxtopdf import convert
# import tr
df = pd.read_excel('book1.xlsx')




#method for shapes update
def update_slide_content(slide, pmarks_text, project_text, viva_text, name_text, msg_text, certNo, dte,total):
    shapes = [slide.shapes[2], slide.shapes[1], slide.shapes[0], slide.shapes[3], slide.shapes[4], slide.shapes[6], slide.shapes[5],slide.shapes[7]]
    text_list = [pmarks_text, project_text, viva_text, name_text, msg_text,certNo, dte, total]

    for shape, text in zip(shapes, text_list):
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = text


# Open the template presentation
templt = Presentation("template/template.pptx")
slide = templt.slides[0]

for i,r in df.iterrows():
    certificateNo="12-12-2024/advetfgvefg"
    dt="12-12-2024"
    updateMsg="of "+ r['Address']+ ", "+ r['Distsrict']+ " has successfully completed the course in "+ r['Course']+" here from "+ r['Session'] +" "+ str(r['Year'])  +" with Grade "+ r['Grade']
    nm=r['Title'] + " " + r['Name'] # title and Name
    viva=str(r['Viva'])
    projectm=str(r['Project'])
    practical=str(r['Practical']) 
    total=str(r['Total']) 
    update_slide_content(slide,viva,projectm,practical,nm,updateMsg, certificateNo, dt,total)
    fname = f"{r['Name']}.pptx"
    oname=f"{r['Name']}"
    oname="OutputFile/" + oname
   
    templt.save("certificate/"+ fname)
    print(fname)
    fname="certificate/" + fname
    convert(fname, oname) #convert to pdf
    
    # tr.covrt(fname)
